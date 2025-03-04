
import streamlit as st
import openai
import time
from collections import deque
from io import BytesIO
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader, PdfWriter
import pandas as pd
from pptx.util import Pt
import os


# Rate limiting parameters
MAX_REQUESTS_PER_MINUTE = 3500
MAX_TOKENS_PER_MINUTE = 90000
WINDOW_SECONDS = 60
requests_timestamps = deque()
tokens_timestamps = deque()
# Đọc API key từ tệp
# API_KEY_FILE = "api_key.txt"
# def get_api_key():
#     if os.path.exists(API_KEY_FILE):
#         with open(API_KEY_FILE, "r") as f:
#             return f.read().strip()
#     return None

# # Khởi tạo API key
# api_key = get_api_key()
# if not api_key:
#     st.error("Không tìm thấy API key. Hãy chắc chắn rằng tệp 'api_key.txt' chứa API key của bạn.")
#     st.stop()
def check_and_wait_for_rate_limit(tokens_used: int):
    current_time = time.time()
    while requests_timestamps and (current_time - requests_timestamps[0] > WINDOW_SECONDS):
        requests_timestamps.popleft()
    while tokens_timestamps and (current_time - tokens_timestamps[0][0] > WINDOW_SECONDS):
        tokens_timestamps.popleft()
    current_requests = len(requests_timestamps)
    current_tokens = sum(t[1] for t in tokens_timestamps)
    if current_requests + 1 > MAX_REQUESTS_PER_MINUTE or current_tokens + tokens_used > MAX_TOKENS_PER_MINUTE:
        time.sleep(1)
        return check_and_wait_for_rate_limit(tokens_used)
    requests_timestamps.append(current_time)
    tokens_timestamps.append((current_time, tokens_used))

def load_specialized_dict_from_excel(excel_file):
    if excel_file is None:
        return {}
    df = pd.read_excel(excel_file)
    return {str(row['English']).strip(): str(row['Vietnamese']).strip() for _, row in df.iterrows() if row['English'] and row['Vietnamese']}

def translate_text_with_chatgpt(original_text, api_key, global_dict=None):
    if not original_text.strip():
        return original_text
    partial_dict = {eng: vie for eng, vie in global_dict.items() if eng.lower() in original_text.lower()} if global_dict else {}
    dict_prompt = "\n".join([f"{k}: {v}" for k, v in partial_dict.items()]) if partial_dict else ""
    user_prompt = f"{dict_prompt}\n\n{original_text}"
    client = openai.OpenAI(api_key=api_key)
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "Bạn là một trợ lý AI dịch thuật. Hãy dịch văn bản sau từ tiếng Anh sang tiếng Việt, ưu tiên dùng đúng các thuật ngữ chuyên ngành (nếu có). Trước tiên hay tra cứu từ vựng trong câu có từ  nào thuộc từ vựng nằm trong file từ vựng chuyên ngành mà tôi cung cấp không, nếu có hãy dùng nghĩa tiếng việt của từ vựng chuyên ngành đó được cung cấp trong file xlsx, các từ còn lại bạn có thể dịch tự động. ** Lưu ý mỗi câu chỉ được phép dịch 1 lần duy nhất, ngoài ra nếu đó là mội chuỗi kí tự bất kì không phải là bất kì từ tiếng anh nào thì đó có thể là kí hiệu hoặc mã của sản phẩm bạn có thể giữ nguyên và không cần dịch sang tiếng việt. Nếu đầu vào (input) không có nội dung thì bạn có thể bỏ qua và không trả về kết quả gì hết ( không trả output)."},
            {"role": "user", "content": user_prompt}
        ],
        temperature=0.2,
        max_tokens=2048
    )
    translated_text = response.choices[0].message.content
    check_and_wait_for_rate_limit(response.usage.total_tokens if response.usage else 0)
    return translated_text


from pptx.util import Pt

def adjust_text_fit(text_frame, shape):
    """
    Adjust text font size dynamically to fit within the text box without overflow.
    Uses shape.width and shape.height instead of text_frame.width.
    """
    max_width = shape.width  # Get the width of the text box
    max_height = shape.height  # Get the height of the text box
    min_font_size = Pt(8)  # Set a minimum font size to maintain readability

    for para in text_frame.paragraphs:
        for run in para.runs:
            if run.font.size and run.font.size > min_font_size:
                run.font.size = max(min_font_size, run.font.size * 0.9)  # Reduce font size if needed

def distribute_text_across_runs(para, translated_text):
    """
    Phân phối đều văn bản đã dịch qua các run trong khi giữ nguyên định dạng.
    """
    original_text = "".join(run.text for run in para.runs)
    if not original_text.strip():
        return

    total_original_len = len(original_text)
    if total_original_len == 0:
        return

    remaining_text = translated_text
    for run in para.runs:
        if not run.text:
            continue

        # Tính toán phân phối văn bản theo tỷ lệ
        run_len = len(run.text)
        portion = min(run_len / total_original_len, 1.0)
        chars_to_take = int(len(translated_text) * portion)

        # Cập nhật văn bản trong khi giữ nguyên định dạng
        run.text = remaining_text[:chars_to_take]
        remaining_text = remaining_text[chars_to_take:]

        # Giữ nguyên định dạng (font, cỡ chữ, in đậm, in nghiêng, màu sắc)
        if run.font is not None:
            run.font.name = run.font.name  # Giữ font chữ
            run.font.size = run.font.size  # Giữ cỡ chữ
            run.font.bold = run.font.bold  # Giữ in đậm
            run.font.italic = run.font.italic  # Giữ in nghiêng
            
            if run.font.color and hasattr(run.font.color, 'rgb'):
                run.font.color.rgb = run.font.color.rgb  # Giữ màu sắc

    # Gắn phần văn bản còn lại vào run cuối cùng (nếu có)
    if remaining_text and para.runs:
        para.runs[-1].text += remaining_text
def delete_unwanted_slides(pr, start_slide, end_slide):
    """
    Xóa các slide không nằm trong khoảng start_slide đến end_slide, đảm bảo hoạt động chính xác ngay cả khi
    chỉ lấy một phạm vi nhỏ từ một tệp lớn.

    Args:
        pr (Presentation): Đối tượng PowerPoint cần xử lý.
        start_slide (int): Slide bắt đầu giữ lại (chỉ mục 1-based).
        end_slide (int): Slide kết thúc giữ lại (chỉ mục 1-based).
    """
    total_slides = len(pr.slides)

    # Kiểm tra nếu phạm vi nằm ngoài số lượng slide có sẵn
    if start_slide < 1 or end_slide > total_slides or start_slide > end_slide:
        raise ValueError("Phạm vi slide không hợp lệ! Vui lòng nhập giá trị hợp lệ.")

    # Danh sách các index slide cần giữ lại
    keep_slides = set(range(start_slide - 1, end_slide))

    # Lấy danh sách các slide ID thực sự trong XML
    xml_slides = pr.slides._sldIdLst
    slides = list(xml_slides)

    # Xóa những slide không thuộc phạm vi cần giữ lại (từ cuối về đầu)
    for i in reversed(range(total_slides)):
        if i not in keep_slides:
            pr.part.drop_rel(slides[i].rId)
            xml_slides.remove(slides[i])

def translate_pptx(pptx_file: BytesIO, api_key: str, specialized_dict: dict[str, str],start_slide: int, end_slide: int) -> BytesIO:
    """
    Dịch văn bản trong file PowerPoint từ tiếng Anh sang tiếng Việt, giữ nguyên font, cỡ chữ và màu sắc gốc.
    Xử lý văn bản tràn bằng cách điều chỉnh cỡ chữ động.
    
    Args:
        pptx_file: Đối tượng BytesIO chứa file PPTX
        api_key: Khóa API OpenAI
        specialized_dict: Từ điển thuật ngữ chuyên ngành (Anh -> Việt)
    
    Returns:
        Đối tượng BytesIO chứa file PPTX đã dịch
    """
    pr = Presentation(pptx_file)
    progress_bar = st.progress(0)
    status_text = st.empty()
    # new_presentation = Presentation()
    # for i, slide in enumerate(pr.slides):
    #     status_text.text(f"Đang dịch slide {i+1}/{total_slides}...")
    # slides_to_delete = [i for i in range(total_slides) if i < start_slide - 1 or i >= end_slide]

    # for index in reversed(slides_to_delete):  # Xóa từ cuối về đầu để tránh lỗi index
    #     xml_slides = pr.slides._sldIdLst  
    #     slides = list(xml_slides)
    #     pr.part.drop_rel(slides[index].rId)
    #     xml_slides.remove(slides[index])
    # delete_unwanted_slides(pr, start_slide, end_slide)
    total_slides = len(pr.slides)
    for i in range(start_slide - 1, end_slide):
    # for i, slide in enumerate(total_slides):
        slide = pr.slides[i]
        status_text.text(f"Đang dịch slide {i+1}/{end_slide}...")

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                for para in text_frame.paragraphs:
                    if not para.text.strip():
                        continue

                    # Thu thập văn bản gốc
                    original_text = "".join(run.text for run in para.runs)
                    translated_text = translate_text_with_chatgpt(original_text, api_key, specialized_dict)

                    # Bỏ qua nếu không có bản dịch hoặc bản dịch không thay đổi
                    if not translated_text or translated_text == original_text or translated_text == 'Xin lỗi, nhưng văn bản bạn cung cấp không đủ để dịch. Bạn có thể cung cấp thêm ngữ cảnh hoặc thông tin chi tiết hơn không?':
                        continue

                    # Phân phối văn bản đã dịch qua các run
                    distribute_text_across_runs(para, translated_text)

                # Điều chỉnh kích thước văn bản để tránh tràn
                adjust_text_fit(text_frame, shape)

            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text_frame = cell.text_frame
                        for para in text_frame.paragraphs:
                            if not para.text.strip():
                                continue
                            # Thu thập văn bản gốc từ các run
                            original_text = "".join(run.text for run in para.runs)
                            translated_text = translate_text_with_chatgpt(original_text, api_key, specialized_dict)
                            # Chỉ cập nhật nếu bản dịch hợp lệ
                            if translated_text and translated_text != original_text and translated_text != 'Xin lỗi, nhưng văn bản bạn cung cấp không đủ để dịch. Bạn có thể cung cấp thêm ngữ cảnh hoặc thông tin chi tiết hơn không?':
                                distribute_text_across_runs(para, translated_text)
                        # Tùy chọn: Điều chỉnh kích thước văn bản trong ô (nếu cần)
                                # adjust_text_fit_for_cell(cell)

        progress_bar.progress((i+1 - (start_slide - 1)) / (end_slide - start_slide + 1))
    delete_unwanted_slides(pr, start_slide, end_slide)
    output = BytesIO()
    pr.save(output)
    output.seek(0)
    status_text.text("Dịch PPTX hoàn tất!")
    return output

# Streamlit UI
st.set_page_config(page_title="Auto Translator App with Full Formatting")
st.title("VN DITECH JSC")
st.subheader("_Tự động dịch tài liệu_ :orange[(PPTX)] . _Giữ nguyên định dạng_", divider="orange")


api_key = st.text_input("Nhập OpenAI API key của bạn:", type="password")
uploaded_excel_dict = st.file_uploader("Tải file từ điển nếu có ( Excel )", type=["xlsx"])
specialized_dict = load_specialized_dict_from_excel(uploaded_excel_dict)

uploaded_file = st.file_uploader("Tải lên file cần dịch (PPTX)", type=["pptx"])

if uploaded_file:
    pr = Presentation(uploaded_file)
    total_slides = len(pr.slides)
    start_slide = st.number_input("Chọn trang bắt đầu", min_value=1, max_value=total_slides, value=1)
    end_slide = st.number_input("Chọn trang kết thúc", min_value=start_slide, max_value=total_slides, value=total_slides)
    
    if st.button("Dịch file PPTX") and api_key:
        output = translate_pptx(uploaded_file, api_key, specialized_dict, start_slide, end_slide)
        file_name = "VN_" + uploaded_file.name
        st.download_button("Tải về file đã dịch", output, file_name )
